# -*- coding: utf-8 -*-
"""
Build the oral defense slide deck for the FR Power thesis.
EDHEC MSc Data Analysis & AI — Cambreleng & Oumedjeber.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

FIG = "thesis/figures"
OUT = "outputs/Defense_Presentation.pptx"

EDHEC = RGBColor(0, 62, 126)
ACCENT = RGBColor(46, 95, 163)
DARK = RGBColor(33, 37, 41)
GREY = RGBColor(110, 110, 110)
LIGHT = RGBColor(238, 243, 249)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(34, 139, 76)
RED = RGBColor(192, 57, 43)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, x, y, w, h, color, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    return sp


def txt(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """lines: list of dicts {text, size, bold, color, space_after, bullet}"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get("align", align)
        if ln.get("space_after") is not None:
            p.space_after = Pt(ln["space_after"])
        if ln.get("space_before") is not None:
            p.space_before = Pt(ln["space_before"])
        runs = ln["text"] if isinstance(ln["text"], list) else [(ln["text"], {})]
        for seg, ov in runs:
            r = p.add_run()
            r.text = seg
            r.font.size = Pt(ov.get("size", ln.get("size", 18)))
            r.font.bold = ov.get("bold", ln.get("bold", False))
            r.font.italic = ov.get("italic", ln.get("italic", False))
            r.font.color.rgb = ov.get("color", ln.get("color", DARK))
            r.font.name = "Calibri"
    return tb


def header(slide, kicker, title, speaker):
    rect(slide, 0, 0, SW, Inches(1.15), EDHEC)
    rect(slide, 0, Inches(1.15), SW, Pt(3), ACCENT)
    txt(slide, Inches(0.55), Inches(0.12), Inches(10.5), Inches(0.4),
        [{"text": kicker.upper(), "size": 12, "bold": True, "color": RGBColor(170, 200, 235)}])
    txt(slide, Inches(0.55), Inches(0.40), Inches(11.6), Inches(0.7),
        [{"text": title, "size": 26, "bold": True, "color": WHITE}])
    # speaker chip
    chip = rect(slide, Inches(11.0), Inches(0.30), Inches(1.9), Inches(0.5), ACCENT)
    txt(slide, Inches(11.0), Inches(0.30), Inches(1.9), Inches(0.5),
        [{"text": speaker, "size": 13, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER}],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def footer(slide, n):
    txt(slide, Inches(0.5), Inches(7.05), Inches(8), Inches(0.35),
        [{"text": "Weather in French Day-Ahead Electricity Price Forecasting", "size": 9, "color": GREY}])
    txt(slide, Inches(12.0), Inches(7.05), Inches(0.9), Inches(0.35),
        [{"text": str(n), "size": 11, "bold": True, "color": EDHEC, "align": PP_ALIGN.RIGHT}], align=PP_ALIGN.RIGHT)


def pic_fit(slide, path, x, y, w, h):
    """Insert image fit into box keeping aspect ratio, centered."""
    from PIL import Image
    iw, ih = Image.open(path).size
    box_ar = w / h
    img_ar = iw / ih
    if img_ar > box_ar:
        nw = w
        nh = int(w / img_ar)
    else:
        nh = h
        nw = int(h * img_ar)
    nx = x + (w - nw) // 2
    ny = y + (h - nh) // 2
    slide.shapes.add_picture(path, nx, ny, nw, nh)


# ============================================================
# SLIDE 1 — TITLE
# ============================================================
s = add_slide()
rect(s, 0, 0, SW, SH, WHITE)
rect(s, 0, 0, SW, Inches(0.28), EDHEC)
rect(s, 0, SH - Inches(0.28), SW, Inches(0.28), EDHEC)
txt(s, Inches(1.0), Inches(0.9), Inches(11.3), Inches(0.5),
    [{"text": "MSc DATA ANALYSIS & AI  ·  EDHEC BUSINESS SCHOOL  ·  MASTER THESIS DEFENSE",
      "size": 13, "bold": True, "color": ACCENT, "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER)
rect(s, Inches(2.0), Inches(1.7), Inches(9.33), Pt(2.5), EDHEC)
txt(s, Inches(1.0), Inches(2.0), Inches(11.3), Inches(2.0),
    [{"text": "The Role of Weather in French Day-Ahead", "size": 34, "bold": True, "color": EDHEC, "align": PP_ALIGN.CENTER},
     {"text": "Electricity Price Forecasting", "size": 34, "bold": True, "color": EDHEC, "align": PP_ALIGN.CENTER},
     {"text": "A Random Forest Approach", "size": 22, "italic": True, "color": ACCENT, "align": PP_ALIGN.CENTER, "space_before": 10}],
    align=PP_ALIGN.CENTER)
rect(s, Inches(2.0), Inches(4.25), Inches(9.33), Pt(2.5), EDHEC)
txt(s, Inches(1.0), Inches(4.5), Inches(11.3), Inches(1.3),
    [{"text": [("Leo Cambreleng", {"bold": True, "size": 20, "color": DARK}),
               ("        Lyam Oumedjeber", {"bold": True, "size": 20, "color": DARK})],
      "align": PP_ALIGN.CENTER},
     {"text": "Supervisor: Prof. Milos Vulanovic", "size": 16, "color": GREY, "align": PP_ALIGN.CENTER, "space_before": 12},
     {"text": "June 2026", "size": 14, "color": GREY, "align": PP_ALIGN.CENTER, "space_before": 6}],
    align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2 — CONTEXT & MOTIVATION (Léo)
# ============================================================
s = add_slide()
header(s, "Context & Motivation", "Why forecast French day-ahead prices?", "LÉO")
bullets = [
    ("Electricity cannot be stored at scale", "supply & demand clear instantaneously every hour → extreme volatility, spikes up to several 000s EUR/MWh"),
    ("France is a special market", "~63 GW nuclear (~70% of output); electric heating → ~2.4 GW/°C winter thermosensitivity — the highest in Europe"),
    ("Forecasts have real economic stakes", "a 1 EUR/MWh accuracy gain on a 100 MW book ≈ 876,000 EUR/year"),
    ("So: does weather actually help?", "temperature & wind drive demand and renewable supply — but is that signal already priced in elsewhere?"),
]
y = Inches(1.55)
for head, sub in bullets:
    rect(s, Inches(0.6), y + Inches(0.05), Inches(0.16), Inches(0.95), ACCENT)
    txt(s, Inches(0.95), y, Inches(11.6), Inches(1.0),
        [{"text": head, "size": 19, "bold": True, "color": EDHEC},
         {"text": sub, "size": 15, "color": DARK, "space_before": 3}])
    y += Inches(1.28)
footer(s, 2)

# ============================================================
# SLIDE 3 — RESEARCH QUESTION (Léo)
# ============================================================
s = add_slide()
header(s, "Research Question", "What we set out to answer", "LÉO")
box = rect(s, Inches(0.7), Inches(1.5), Inches(11.9), Inches(1.5), LIGHT)
txt(s, Inches(1.0), Inches(1.6), Inches(11.3), Inches(1.3),
    [{"text": [("“Do weather variables significantly improve the accuracy of French day-ahead "
                "electricity price forecasts — and does their contribution depend on the market regime?”",
                {"size": 19, "italic": True, "bold": True, "color": EDHEC})]}],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
subs = [
    ("1", "Can ML (Random Forest, XGBoost) beat a naïve statistical benchmark?"),
    ("2", "Do weather features add statistically significant value over fundamentals alone?"),
    ("3", "Does that value differ between a stable regime (2024–25) and the 2022 crisis — and why?"),
    ("4", "Can the forecasts be turned into real trading profit under realistic costs?"),
]
y = Inches(3.35)
for num, q in subs:
    c = rect(s, Inches(0.7), y, Inches(0.55), Inches(0.7), EDHEC)
    txt(s, Inches(0.7), y, Inches(0.55), Inches(0.7),
        [{"text": num, "size": 22, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER}],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(1.45), y, Inches(11.2), Inches(0.7),
        [{"text": q, "size": 16, "color": DARK}], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.88)
footer(s, 3)

# ============================================================
# SLIDE 4 — DATA (Léo)
# ============================================================
s = add_slide()
header(s, "Data", "A reproducible French dataset, 2018–2025", "LÉO")
txt(s, Inches(0.55), Inches(1.35), Inches(6.0), Inches(0.5),
    [{"text": "~64,000 hourly observations  ·  4 public sources", "size": 15, "bold": True, "color": ACCENT}])
cards = [
    ("ENTSO-E", "Prices, load forecast, generation by source, cross-border flows"),
    ("ERA5 (ECMWF)", "Temperature, wind, solar irradiance, precipitation — spatial mean France"),
    ("Fuel prices", "TTF natural gas (EUR/MWh), ARA coal (EUR/t)"),
    ("Engineered", "35 features: calendar, price lags, fundamentals, flows, weather, derived"),
]
y = Inches(1.95)
for i, (t, d) in enumerate(cards):
    cy = y + Inches(i * 1.18)
    rect(s, Inches(0.55), cy, Inches(6.1), Inches(1.02), LIGHT)
    rect(s, Inches(0.55), cy, Inches(0.14), Inches(1.02), ACCENT)
    txt(s, Inches(0.85), cy + Inches(0.08), Inches(5.7), Inches(0.9),
        [{"text": t, "size": 16, "bold": True, "color": EDHEC},
         {"text": d, "size": 12.5, "color": DARK, "space_before": 2}])
pic_fit(s, f"{FIG}/price_series_full.png", Inches(6.9), Inches(1.65), Inches(6.0), Inches(4.0))
txt(s, Inches(6.9), Inches(5.7), Inches(6.0), Inches(0.6),
    [{"text": "Three regimes: pre-crisis (2018–21), the 2022 energy crisis (>1,000 EUR/MWh), and normalisation (2023+).",
      "size": 11.5, "italic": True, "color": GREY, "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER)
footer(s, 4)

# ============================================================
# SLIDE 5 — THE FOUR MODELS (Léo)
# ============================================================
s = add_slide()
header(s, "Methodology", "A nested four-model ablation design", "LÉO")
models = [
    ("A", "Naïve benchmark", "price 168h ago (same hour, last week)", "The bar to beat", GREY),
    ("B", "RF without weather", "27 features — no meteorology", "Ablation baseline", ACCENT),
    ("C", "RF with weather", "35 features — full meteorology", "Main model", EDHEC),
    ("D", "XGBoost with weather", "same 35 features, boosted trees", "More expressive?", RGBColor(180, 95, 30)),
]
y = Inches(1.5)
for letter, name, desc, role, col in models:
    rect(s, Inches(0.7), y, Inches(1.0), Inches(1.05), col)
    txt(s, Inches(0.7), y, Inches(1.0), Inches(1.05),
        [{"text": letter, "size": 36, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER}],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, Inches(1.85), y, Inches(10.8), Inches(1.05), LIGHT)
    txt(s, Inches(2.1), y + Inches(0.1), Inches(8.0), Inches(0.9),
        [{"text": name, "size": 18, "bold": True, "color": col},
         {"text": desc, "size": 14, "color": DARK, "space_before": 2}])
    txt(s, Inches(10.0), y, Inches(2.4), Inches(1.05),
        [{"text": role, "size": 13, "italic": True, "bold": True, "color": col, "align": PP_ALIGN.RIGHT}],
        align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(1.22)
txt(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.5),
    [{"text": [("Key comparison:  ", {"bold": True, "color": EDHEC, "size": 15}),
               ("C vs B isolates the marginal value of weather.  Evaluated out-of-sample with the Diebold-Mariano test (HLN-corrected).",
                {"size": 15, "color": DARK})]}])
footer(s, 5)

# ============================================================
# SLIDE 6 — MAIN RESULTS 2024-25 (Léo)
# ============================================================
s = add_slide()
header(s, "Results — Stable Regime (May 2024–Apr 2025)", "Machine learning works; weather barely moves the needle", "LÉO")
# table
rows = [
    ("Model", "MAE", "RMSE", "R²", "Hit %"),
    ("A — Naïve", "33.09", "44.23", "0.160", "59.6"),
    ("B — RF no weather", "16.95", "22.88", "0.775", "65.5"),
    ("C — RF weather", "16.94", "22.85", "0.776", "65.4"),
    ("D — XGBoost", "19.14", "28.18", "0.659", "65.8"),
]
tx, ty, tw = Inches(0.55), Inches(1.55), Inches(6.4)
rh = Inches(0.62)
colw = [Inches(2.4), Inches(1.0), Inches(1.0), Inches(1.0), Inches(1.0)]
for ri, row in enumerate(rows):
    cy = ty + rh * ri
    if ri == 0:
        rect(s, tx, cy, tw, rh, EDHEC)
    elif row[0].startswith("C"):
        rect(s, tx, cy, tw, rh, RGBColor(214, 232, 240))
    else:
        rect(s, tx, cy, tw, rh, LIGHT if ri % 2 else WHITE)
    cx = tx
    for ci, cell in enumerate(row):
        col = WHITE if ri == 0 else (EDHEC if row[0].startswith("C") else DARK)
        bold = ri == 0 or row[0].startswith("C") or ci == 0
        al = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
        txt(s, cx + Inches(0.1), cy, colw[ci] - Inches(0.1), rh,
            [{"text": cell, "size": 14, "bold": bold, "color": col, "align": al}],
            align=al, anchor=MSO_ANCHOR.MIDDLE)
        cx += colw[ci]
txt(s, Inches(0.55), Inches(4.85), Inches(6.4), Inches(2.0),
    [{"text": [("− 49% MAE", {"bold": True, "size": 17, "color": GREEN}),
               (" vs naïve (33.1 → 16.9 EUR/MWh), R² 0.16 → 0.78", {"size": 15, "color": DARK})],
      "space_after": 8},
     {"text": [("Weather C vs B: ", {"bold": True, "size": 15, "color": EDHEC}),
               ("MAE −20.01 EUR/MWh — negligible", {"size": 15, "color": DARK})], "space_after": 8},
     {"text": [("XGBoost underperforms RF here", {"size": 15, "color": DARK})]}])
pic_fit(s, f"{FIG}/model_comparison.png", Inches(7.1), Inches(1.5), Inches(5.9), Inches(5.0))
footer(s, 6)

# ============================================================
# SLIDE 7 — CENTRAL FINDING / DM TEST (Léo -> handoff)
# ============================================================
s = add_slide()
header(s, "Central Finding", "Weather is not statistically significant — in the stable regime", "LÉO")
box = rect(s, Inches(0.7), Inches(1.45), Inches(11.9), Inches(1.15), RGBColor(214, 232, 240))
txt(s, Inches(1.0), Inches(1.5), Inches(11.3), Inches(1.05),
    [{"text": [("Diebold-Mariano  C vs B:   ", {"bold": True, "size": 19, "color": EDHEC}),
               ("p = 0.572   (DM = −0.565)   →   fail to reject equal accuracy",
                {"size": 19, "bold": True, "color": RED})]}],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rows = [
    ("Comparison", "DM stat", "p-value", "Sig."),
    ("C vs B — does weather add value?", "−0.565", "0.572", "n.s."),
    ("C vs A — RF-weather beats naïve?", "−53.53", "<0.001", "***"),
    ("B vs A — RF-no-weather beats naïve?", "−53.44", "<0.001", "***"),
    ("D vs C — XGBoost vs RF-weather?", "+10.88", "<0.001", "***"),
]
tx, ty, tw = Inches(0.7), Inches(2.95), Inches(11.9)
rh = Inches(0.6)
colw = [Inches(6.5), Inches(1.8), Inches(1.8), Inches(1.8)]
for ri, row in enumerate(rows):
    cy = ty + rh * ri
    if ri == 0:
        rect(s, tx, cy, tw, rh, EDHEC)
    elif ri == 1:
        rect(s, tx, cy, tw, rh, RGBColor(250, 224, 220))
    else:
        rect(s, tx, cy, tw, rh, LIGHT if ri % 2 else WHITE)
    cx = tx
    for ci, cell in enumerate(row):
        col = WHITE if ri == 0 else (RED if ri == 1 else DARK)
        al = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
        txt(s, cx + Inches(0.15), cy, colw[ci] - Inches(0.15), rh,
            [{"text": cell, "size": 14, "bold": (ri == 0 or ri == 1 or ci == 0), "color": col, "align": al}],
            align=al, anchor=MSO_ANCHOR.MIDDLE)
        cx += colw[ci]
txt(s, Inches(0.7), Inches(6.15), Inches(11.9), Inches(0.7),
    [{"text": [("ML adds robust value; weather does not — yet.  ", {"bold": True, "size": 15, "color": EDHEC}),
               ("Why? And is this always true? → over to Lyam.", {"italic": True, "size": 15, "color": GREY})]}])
footer(s, 7)

# ============================================================
# SLIDE 8 — FEATURE IMPORTANCE (Lyam)
# ============================================================
s = add_slide()
header(s, "Feature Importance", "Where does the predictive power come from?", "LYAM")
pic_fit(s, f"{FIG}/feature_importance_rf.png", Inches(0.5), Inches(1.45), Inches(7.3), Inches(5.3))
txt(s, Inches(8.0), Inches(1.55), Inches(4.9), Inches(5.0),
    [{"text": "Price lags dominate", "size": 18, "bold": True, "color": EDHEC, "space_after": 4},
     {"text": "lag-24h alone = 20.7% of importance; lags + rolling means ≈ 70%", "size": 14, "color": DARK, "space_after": 12},
     {"text": "Fuel prices matter", "size": 18, "bold": True, "color": EDHEC, "space_after": 4},
     {"text": "TTF gas 8.3%, ARA coal 6.0%, nuclear availability ~1.9%", "size": 14, "color": DARK, "space_after": 12},
     {"text": "Weather is small", "size": 18, "bold": True, "color": RED, "space_after": 4},
     {"text": "all weather features combined ≈ 2.2% of total importance", "size": 14, "color": DARK, "space_after": 12},
     {"text": [("→ consistent with the DM test: weather signal is already captured by load forecast & TTF.",
                {"size": 14, "italic": True, "color": ACCENT})]}])
footer(s, 8)

# ============================================================
# SLIDE 9 — INFORMATION REDUNDANCY HYPOTHESIS (Lyam)
# ============================================================
s = add_slide()
header(s, "Interpretation", "The Information Redundancy Hypothesis", "LYAM")
txt(s, Inches(0.7), Inches(1.4), Inches(11.9), Inches(0.7),
    [{"text": "Weather looks redundant because its signal is already encoded in other features:",
      "size": 17, "bold": True, "color": EDHEC}])
chans = [
    ("Load forecast ≡ temperature", "RTE's published load forecast is built on temperature forecasts — including it implicitly conditions on weather-driven demand."),
    ("TTF gas ≡ European weather", "Gas prices rise with pan-European heating demand — the commodity market aggregates the weather signal."),
    ("Nuclear dominance", "The nuclear availability ratio explains a large share of price variance independently of weather."),
]
y = Inches(2.25)
for t, d in chans:
    rect(s, Inches(0.7), y, Inches(11.9), Inches(1.25), LIGHT)
    rect(s, Inches(0.7), y, Inches(0.16), Inches(1.25), ACCENT)
    txt(s, Inches(1.05), y + Inches(0.12), Inches(11.3), Inches(1.05),
        [{"text": t, "size": 17, "bold": True, "color": EDHEC},
         {"text": d, "size": 14.5, "color": DARK, "space_before": 3}])
    y += Inches(1.45)
footer(s, 9)

# ============================================================
# SLIDE 10 — 2022 CRISIS REVERSAL (Lyam)
# ============================================================
s = add_slide()
header(s, "Robustness — 2022 Energy Crisis", "The result reverses: weather becomes highly significant", "LYAM")
txt(s, Inches(0.55), Inches(1.35), Inches(6.3), Inches(0.6),
    [{"text": "Models retrained on 2018–2021, tested on full-year 2022 (strict temporal split).",
      "size": 13.5, "italic": True, "color": GREY}])
# comparison cards
def stat_card(slide, x, y, w, title, p, dm, col):
    rect(slide, x, y, w, Inches(1.55), col)
    txt(slide, x, y + Inches(0.12), w, Inches(1.4),
        [{"text": title, "size": 15, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER},
         {"text": p, "size": 26, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER, "space_before": 6},
         {"text": dm, "size": 13, "color": RGBColor(225, 235, 245), "align": PP_ALIGN.CENTER, "space_before": 2}],
        align=PP_ALIGN.CENTER)

stat_card(s, Inches(0.55), Inches(2.0), Inches(3.05), "Stable 2024–25", "p = 0.572", "DM = −0.565  ·  n.s.", GREY)
stat_card(s, Inches(3.75), Inches(2.0), Inches(3.05), "Crisis 2022", "p < 0.001", "DM = −13.27  ·  ***", RED)
txt(s, Inches(0.55), Inches(3.8), Inches(6.3), Inches(2.6),
    [{"text": "What happens in 2022:", "size": 16, "bold": True, "color": EDHEC, "space_after": 6},
     {"text": "All models degrade sharply — MAE 16.9 → 66.6 EUR/MWh (~4×)", "size": 14, "color": DARK, "space_after": 6},
     {"text": "But C now beats B by 0.73 EUR/MWh, highly significant", "size": 14, "color": DARK, "space_after": 6},
     {"text": "XGBoost degrades most (+14.6% vs RF) — RF is the robust choice", "size": 14, "color": DARK}])
pic_fit(s, f"{FIG}/model_comparison_2022_vs_2024.png", Inches(7.0), Inches(1.5), Inches(6.0), Inches(5.1))
footer(s, 10)

# ============================================================
# SLIDE 11 — WHY THE CRISIS BREAKS REDUNDANCY (Lyam)
# ============================================================
s = add_slide()
header(s, "Interpretation", "Why the crisis breaks the redundancy", "LYAM")
txt(s, Inches(0.7), Inches(1.4), Inches(11.9), Inches(0.6),
    [{"text": "The same three channels that made weather redundant break down in 2022:",
      "size": 17, "bold": True, "color": EDHEC}])
mech = [
    ("TTF decouples from local weather", "The Russian gas supply shock drove TTF by geopolitics, not French temperature — the gas–weather link snaps."),
    ("Nuclear constraints amplify demand", "With availability ~40%, there is no headroom to absorb cold-driven demand surges → temperature feeds straight into price."),
    ("Signal-to-noise rises", "When prices swing by hundreds of EUR/MWh, the direct temperature–demand effect becomes large enough to measure."),
]
y = Inches(2.2)
for i, (t, d) in enumerate(mech, 1):
    rect(s, Inches(0.7), y, Inches(0.6), Inches(1.2), RED)
    txt(s, Inches(0.7), y, Inches(0.6), Inches(1.2),
        [{"text": str(i), "size": 24, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER}],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, Inches(1.45), y, Inches(11.15), Inches(1.2), LIGHT)
    txt(s, Inches(1.75), y + Inches(0.12), Inches(10.7), Inches(1.0),
        [{"text": t, "size": 16.5, "bold": True, "color": EDHEC},
         {"text": d, "size": 14, "color": DARK, "space_before": 3}])
    y += Inches(1.4)
txt(s, Inches(0.7), Inches(6.55), Inches(11.9), Inches(0.5),
    [{"text": [("Takeaway:  ", {"bold": True, "size": 15, "color": EDHEC}),
               ("the value of weather is regime-dependent — a dynamic, not a permanent, modelling choice.",
                {"italic": True, "size": 15, "color": DARK})]}])
footer(s, 11)

# ============================================================
# SLIDE 12 — TRADING BACKTEST (Lyam)
# ============================================================
s = add_slide()
header(s, "Economic Value — Trading Backtest", "Do better forecasts make money?", "LYAM")
txt(s, Inches(0.55), Inches(1.35), Inches(6.3), Inches(0.55),
    [{"text": "Day-ahead directional strategy, central cost 0.30 EUR/MWh, 1 MW position.",
      "size": 13, "italic": True, "color": GREY}])
rows = [
    ("Model", "Net P&L", "Sharpe", "MaxDD", "Calmar"),
    ("A — Naïve", "85,482", "10.5", "2,367", "37"),
    ("B — RF no wx", "136,884", "19.4", "433", "321"),
    ("C — RF wx", "136,711", "19.5", "422", "328"),
    ("D — XGBoost", "142,012", "18.8", "1,351", "107"),
]
tx, ty, tw = Inches(0.55), Inches(2.0), Inches(6.3)
rh = Inches(0.6)
colw = [Inches(2.0), Inches(1.4), Inches(0.95), Inches(1.05), Inches(0.9)]
for ri, row in enumerate(rows):
    cy = ty + rh * ri
    if ri == 0:
        rect(s, tx, cy, tw, rh, EDHEC)
    elif row[0].startswith("C"):
        rect(s, tx, cy, tw, rh, RGBColor(214, 232, 240))
    else:
        rect(s, tx, cy, tw, rh, LIGHT if ri % 2 else WHITE)
    cx = tx
    for ci, cell in enumerate(row):
        col = WHITE if ri == 0 else (EDHEC if row[0].startswith("C") else DARK)
        al = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
        txt(s, cx + Inches(0.08), cy, colw[ci] - Inches(0.08), rh,
            [{"text": cell, "size": 13, "bold": (ri == 0 or row[0].startswith("C") or ci == 0), "color": col, "align": al}],
            align=al, anchor=MSO_ANCHOR.MIDDLE)
        cx += colw[ci]
txt(s, Inches(0.55), Inches(5.2), Inches(6.3), Inches(1.6),
    [{"text": [("RF ≈ 137k EUR/MW", {"bold": True, "size": 15, "color": GREEN}),
               (" vs 85k naïve; long-only ≈ 0 → real skill, not trend", {"size": 13.5, "color": DARK})], "space_after": 6},
     {"text": [("RF drawdown 5.6× smaller", {"bold": True, "size": 15, "color": GREEN}),
               (" than naïve; positive Sharpe every single month", {"size": 13.5, "color": DARK})]}])
pic_fit(s, f"{FIG}/equity_curves_net.png", Inches(7.0), Inches(1.5), Inches(6.0), Inches(5.1))
footer(s, 12)

# ============================================================
# SLIDE 13 — ROBUSTNESS (Lyam)
# ============================================================
s = add_slide()
header(s, "Robustness Checks", "The result holds under different assumptions", "LYAM")
txt(s, Inches(0.55), Inches(1.35), Inches(5.8), Inches(0.5),
    [{"text": "Cost sensitivity — Net P&L (EUR/MW)", "size": 14, "bold": True, "color": EDHEC}])
cost_rows = [
    ("Model", "0.10 EUR/MWh", "0.30 EUR/MWh", "0.60 EUR/MWh"),
    ("A — Naïve", "88,340", "85,482", "81,195"),
    ("B — RF no wx", "139,742", "136,884", "132,597"),
    ("C — RF wx", "139,569", "136,711", "132,424"),
    ("D — XGBoost", "144,870", "142,012", "137,725"),
]
tx2, ty2, tw2 = Inches(0.55), Inches(1.95), Inches(5.8)
rh2 = Inches(0.56)
colw2 = [Inches(1.85), Inches(1.32), Inches(1.32), Inches(1.31)]
for ri, row in enumerate(cost_rows):
    cy = ty2 + rh2 * ri
    if ri == 0:
        rect(s, tx2, cy, tw2, rh2, EDHEC)
    elif row[0].startswith("C"):
        rect(s, tx2, cy, tw2, rh2, RGBColor(214, 232, 240))
    else:
        rect(s, tx2, cy, tw2, rh2, LIGHT if ri % 2 else WHITE)
    cx = tx2
    for ci, cell in enumerate(row):
        col = WHITE if ri == 0 else (EDHEC if row[0].startswith("C") else DARK)
        al = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
        txt(s, cx + Inches(0.07), cy, colw2[ci] - Inches(0.07), rh2,
            [{"text": cell, "size": 12, "bold": (ri == 0 or row[0].startswith("C") or ci == 0), "color": col, "align": al}],
            align=al, anchor=MSO_ANCHOR.MIDDLE)
        cx += colw2[ci]
txt(s, Inches(0.55), Inches(5.1), Inches(5.8), Inches(1.6),
    [{"text": [("RF Sharpe: 19.51 → 19.46 → 19.16", {"bold": True, "size": 14, "color": GREEN}),
               ("  across all three cost scenarios", {"size": 13, "color": DARK})], "space_after": 6},
     {"text": [("Positive Sharpe in ", {"size": 13, "color": DARK}),
               ("all 12 months", {"bold": True, "size": 13, "color": EDHEC}),
               (" of the test period", {"size": 13, "color": DARK})], "space_after": 6},
     {"text": [("DM result stable: ", {"bold": True, "size": 13, "color": EDHEC}),
               ("p = 0.572 n.s. confirmed across multiple sub-periods", {"size": 13, "color": DARK})]}])
pic_fit(s, f"{FIG}/cost_sensitivity.png", Inches(6.8), Inches(1.45), Inches(6.15), Inches(5.1))
footer(s, 13)

# ============================================================
# SLIDE 14 — LIMITATIONS (Lyam)
# ============================================================
s = add_slide()
header(s, "Limitations & Future Work", "What we would do differently", "LYAM")
limits = [
    ("ERA5 = upper bound", "Reanalysis gives perfect hindsight weather — production would use NWP forecasts with 10–30% error. Our result is a ceiling on weather value."),
    ("No walk-forward retraining", "Single batch training; monthly retraining on an expanding window would better reflect real trading conditions."),
    ("Carbon prices (EUA) excluded", "EUA correlates with TTF/coal — partly absorbed — but should be tested explicitly as an additional feature."),
    ("XGBoost not fully tuned", "Default hyperparameters; Bayesian optimisation could close or reverse the RF vs XGBoost gap."),
    ("France only", "The redundancy mechanism is specific to France's nuclear-dominant, high-thermosensitivity structure. Germany, Spain or UK may differ substantially."),
    ("No market impact", "Transaction costs modelled as flat fee; large positions (>10 MW) would face liquidity constraints not captured here."),
]
col_x = [Inches(0.55), Inches(6.85)]
col_w = Inches(6.0)
for idx, (t, d) in enumerate(limits):
    col = idx % 2
    row = idx // 2
    x = col_x[col]
    y = Inches(1.5) + Inches(1.6) * row
    rect(s, x, y, Inches(0.55), Inches(1.35), RGBColor(180, 95, 30))
    txt(s, x, y, Inches(0.55), Inches(1.35),
        [{"text": str(idx + 1), "size": 20, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER}],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, x + Inches(0.65), y, col_w - Inches(0.65), Inches(1.35), LIGHT)
    txt(s, x + Inches(0.75), y + Inches(0.08), col_w - Inches(0.8), Inches(1.2),
        [{"text": t, "size": 14.5, "bold": True, "color": EDHEC},
         {"text": d, "size": 12, "color": DARK, "space_before": 2}])
footer(s, 14)

# ============================================================
# SLIDE 15 — CONCLUSIONS (Lyam)
# ============================================================
s = add_slide()
header(s, "Conclusions", "Four takeaways", "LYAM")
concl = [
    ("ML beats naïve, decisively", "−49% MAE, R² 0.78, p < 0.001 — robust across all 12 months."),
    ("Weather is regime-dependent", "n.s. in the stable regime (redundancy), but significant in the 2022 crisis (redundancy breaks). The core contribution."),
    ("Random Forest > XGBoost", "RF wins on risk-adjusted metrics (Calmar 328 vs 107) and degrades less under stress."),
    ("Forecasts create real value", "~137k EUR/MW net, Sharpe ~19.5, tiny drawdown — robust to transaction costs."),
]
y = Inches(1.5)
for i, (t, d) in enumerate(concl, 1):
    rect(s, Inches(0.7), y, Inches(0.6), Inches(1.15), EDHEC)
    txt(s, Inches(0.7), y, Inches(0.6), Inches(1.15),
        [{"text": str(i), "size": 24, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER}],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, Inches(1.45), y, Inches(11.15), Inches(1.15), LIGHT)
    txt(s, Inches(1.75), y + Inches(0.1), Inches(10.7), Inches(1.0),
        [{"text": t, "size": 17, "bold": True, "color": EDHEC},
         {"text": d, "size": 14, "color": DARK, "space_before": 2}])
    y += Inches(1.34)
footer(s, 15)

# ============================================================
# SLIDE 16 — THANK YOU / Q&A
# ============================================================
s = add_slide()
rect(s, 0, 0, SW, SH, EDHEC)
rect(s, Inches(2.0), Inches(2.6), Inches(9.33), Pt(2.5), ACCENT)
txt(s, Inches(1.0), Inches(2.9), Inches(11.3), Inches(1.4),
    [{"text": "Thank you", "size": 48, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER},
     {"text": "Questions & Discussion", "size": 22, "italic": True, "color": RGBColor(180, 205, 235), "align": PP_ALIGN.CENTER, "space_before": 10}],
    align=PP_ALIGN.CENTER)
rect(s, Inches(2.0), Inches(4.7), Inches(9.33), Pt(2.5), ACCENT)
txt(s, Inches(1.0), Inches(5.0), Inches(11.3), Inches(1.0),
    [{"text": "Leo Cambreleng  ·  Lyam Oumedjeber", "size": 18, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER},
     {"text": "Supervisor: Prof. Milos Vulanovic  ·  EDHEC MSc Data Analysis & AI", "size": 14, "color": RGBColor(180, 205, 235), "align": PP_ALIGN.CENTER, "space_before": 6}],
    align=PP_ALIGN.CENTER)

os.makedirs("outputs", exist_ok=True)
prs.save(OUT)
import sys
print("Saved", OUT, "with", len(prs.slides._sldIdLst), "slides", file=sys.stderr)
