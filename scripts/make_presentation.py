"""
Generate EDHEC thesis defense PowerPoint — accessible, clean version.
Run: python scripts/make_presentation.py
Output: outputs/Soutenance_FR_Power_Prices.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
from pathlib import Path

# ── Colours ─────────────────────────────────────────────────────
NAVY  = RGBColor(0x00, 0x22, 0x55)
GOLD  = RGBColor(0xC8, 0xA0, 0x32)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGREY = RGBColor(0xF4, 0xF6, 0xF9)
DGREY = RGBColor(0x1A, 0x1A, 0x2E)
MGREY = RGBColor(0x6B, 0x7A, 0x8D)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
RED   = RGBColor(0xC0, 0x39, 0x2B)
BLUE  = RGBColor(0x21, 0x6B, 0xB4)

SW = Inches(13.33)
SH = Inches(7.50)

OUT = Path(r"C:\Users\Public\fr-power-thesis\outputs\Soutenance_FR_Power_Prices.pptx")
OUT.parent.mkdir(parents=True, exist_ok=True)

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

HEADER_H  = Inches(1.35)
ACCENT_H  = Inches(0.06)
FOOTER_Y  = Inches(7.20)
FOOTER_H  = Inches(0.30)
CY        = HEADER_H + ACCENT_H + Inches(0.12)   # content top


# ════════════════════════════════════════════════════════════════
# PRIMITIVES
# ════════════════════════════════════════════════════════════════

def emu(v):
    return Emu(max(9144, int(v)))

def rect(sl, x, y, w, h, fill=None):
    w = emu(w); h = emu(h)
    s = sl.shapes.add_shape(1, int(x), int(y), int(w), int(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    s.line.fill.background()
    return s

def txb(sl, text, x, y, w, h,
        size=16, bold=False, italic=False,
        color=DGREY, align=PP_ALIGN.LEFT,
        font="Calibri", wrap=True):
    w = emu(w); h = emu(h)
    tb = sl.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text; run.font.size = Pt(size)
    run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color; run.font.name = font
    return tb

def header(sl, title, subtitle=None):
    rect(sl, 0, 0, SW, HEADER_H, NAVY)
    rect(sl, 0, HEADER_H, SW, ACCENT_H, GOLD)
    txb(sl, title, Inches(0.45), Inches(0.18), Inches(12.4), Inches(0.72),
        size=23, bold=True, color=WHITE)
    if subtitle:
        txb(sl, subtitle, Inches(0.45), Inches(0.90), Inches(12.4), Inches(0.38),
            size=12, color=GOLD)

def footer(sl, n, total=17):
    rect(sl, 0, FOOTER_Y, SW, FOOTER_H, NAVY)
    txb(sl, "Oumedjeber & Cambreleng  |  MSc DAAI  |  EDHEC Business School  |  2025-2026",
        Inches(0.20), FOOTER_Y, Inches(11.5), FOOTER_H, size=9, color=WHITE)
    txb(sl, f"{n} / {total}", Inches(12.0), FOOTER_Y, Inches(1.1), FOOTER_H,
        size=9, color=WHITE, align=PP_ALIGN.RIGHT)

def bg(sl):
    rect(sl, 0, 0, SW, SH, LGREY)

def dot(sl, text, x, y, w, size=13, dc=GOLD, tc=DGREY):
    rect(sl, x, y + Inches(0.14), Inches(0.10), Inches(0.10), dc)
    txb(sl, text, x + Inches(0.18), y, max(Inches(0.5), w - Inches(0.20)),
        Inches(0.45), size=size, color=tc)

def kpi(sl, x, y, w, h, value, label, vc=NAVY, bc=WHITE):
    rect(sl, x, y, w, h, bc)
    rect(sl, x, y, Inches(0.07), h, GOLD)
    txb(sl, value, x+Inches(0.14), y+Inches(0.06),
        max(Inches(0.3), w-Inches(0.18)), Inches(0.52),
        size=26, bold=True, color=vc)
    txb(sl, label, x+Inches(0.14), y+Inches(0.56),
        max(Inches(0.3), w-Inches(0.18)), Inches(0.32),
        size=10, color=MGREY)

def pptx_table(sl, rows, col_widths, x, y, row_h=Inches(0.46),
               hfill=NAVY, hcolor=WHITE, alt=LGREY, body=DGREY,
               hsize=11, bsize=11, hilite=None, hifill=None,
               col_align=None):
    nr = len(rows); nc = len(col_widths)
    tw = sum(col_widths)
    tbl = sl.shapes.add_table(nr, nc, int(x), int(y),
                               int(tw), int(row_h*nr)).table
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = int(cw)
    for ri, row_data in enumerate(rows):
        tbl.rows[ri].height = int(row_h)
        is_h = (ri == 0)
        for ci, val in enumerate(row_data):
            cell = tbl.cell(ri, ci)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            if col_align:
                p.alignment = col_align[ci]
            else:
                p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size  = Pt(hsize if is_h else bsize)
            run.font.bold  = is_h or (hilite is not None and ri-1 == hilite)
            run.font.name  = "Calibri"
            run.font.color.rgb = hcolor if is_h else body
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            sf = etree.SubElement(tcPr, qn('a:solidFill'))
            if is_h:
                clr = hfill
            elif hilite is not None and ri-1 == hilite and hifill:
                clr = hifill
            elif ri % 2 == 1:
                clr = WHITE
            else:
                clr = alt
            sc = etree.SubElement(sf, qn('a:srgbClr'))
            sc.set('val', f'{int(clr[0]):02X}{int(clr[1]):02X}{int(clr[2]):02X}')
            for bt in ('a:lnL','a:lnR','a:lnT','a:lnB'):
                b = etree.SubElement(tcPr, qn(bt))
                etree.SubElement(b, qn('a:noFill'))
    return tbl


# ════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════
def s01():
    sl = prs.slides.add_slide(BLANK)
    rect(sl, 0, 0, SW, SH, NAVY)
    rect(sl, 0, 0, Inches(0.22), SH, GOLD)
    rect(sl, Inches(0.22), Inches(4.30), SW, Inches(0.06), GOLD)
    txb(sl, "FORECASTING DAY-AHEAD ELECTRICITY PRICES",
        Inches(0.55), Inches(1.10), Inches(12.0), Inches(0.52), size=12, color=GOLD)
    txb(sl, "IN FRANCE:\nWHEN DOES WEATHER MATTER?",
        Inches(0.55), Inches(1.68), Inches(11.5), Inches(1.90), size=34, bold=True, color=WHITE)
    txb(sl, "Master Thesis Defense  —  MSc Data Analytics & Artificial Intelligence",
        Inches(0.55), Inches(4.50), Inches(11.0), Inches(0.40), size=13, color=WHITE)
    txb(sl, "Lyam Oumedjeber  &  Leo Cambreleng",
        Inches(0.55), Inches(4.98), Inches(11.0), Inches(0.42), size=16, bold=True, color=WHITE)
    txb(sl, "Thesis Director: Prof. Milos Vulanovic",
        Inches(0.55), Inches(5.48), Inches(11.0), Inches(0.38), size=13, color=GOLD)
    txb(sl, "EDHEC Business School  |  2025-2026",
        Inches(0.55), Inches(5.96), Inches(11.0), Inches(0.38), size=12, color=MGREY)


# ════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ════════════════════════════════════════════════════════════════
def s02():
    sl = prs.slides.add_slide(BLANK)
    bg(sl); header(sl, "Agenda", "20-minute presentation · Questions & discussion to follow")
    footer(sl, 2)

    sections = [
        ("01", "Why this topic? — The French electricity market",     "1 min",  NAVY),
        ("02", "Our data — What we measured and how",                 "2 min",  BLUE),
        ("03", "Four forecasting models — How they were built",       "3 min",  NAVY),
        ("04", "Results — The weather paradox",                       "5 min",  RED),
        ("05", "The 2022 energy crisis — When weather changed everything", "4 min", NAVY),
        ("06", "Trading test — Does better forecasting make money?",  "3 min",  GREEN),
        ("07", "Conclusions",                                         "2 min",  NAVY),
    ]

    ROW_H = Inches(0.70); y0 = CY + Inches(0.05)
    for i, (num, title, dur, col) in enumerate(sections):
        y = y0 + i * ROW_H
        rect(sl, Inches(0.40), y+Inches(0.08), Inches(0.50), Inches(0.50), col)
        txb(sl, num, Inches(0.40), y+Inches(0.08), Inches(0.50), Inches(0.50),
            size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txb(sl, title, Inches(1.08), y+Inches(0.10), Inches(10.6), Inches(0.48),
            size=17, bold=(i in (3,4)), color=DGREY)
        txb(sl, dur, Inches(11.9), y+Inches(0.10), Inches(1.10), Inches(0.48),
            size=12, color=MGREY, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════
# SLIDE 3 — CONTEXT & QUESTION
# ════════════════════════════════════════════════════════════════
def s03():
    sl = prs.slides.add_slide(BLANK)
    bg(sl)
    header(sl, "01 — Why This Topic?", "The French electricity market, briefly explained")
    footer(sl, 3)

    txb(sl, "What is the day-ahead electricity market?",
        Inches(0.40), CY, Inches(5.90), Inches(0.38), size=15, bold=True, color=NAVY)

    items_l = [
        "Every day at noon, energy producers and buyers\nsubmit bids for each hour of the next day",
        "A single price is set for each of the 24 hours — the clearing price",
        "France's market (EPEX SPOT) is the 2nd largest\nin Europe: €40 billion traded per year",
        "Prices are highly volatile: from −€100 to +€4,000/MWh\n(especially during the 2022 energy crisis)",
        "France is unique: ~70% nuclear — very different\ndynamics from Germany or the UK",
        "Academic research on France is scarce —\nmost studies focus on Germany or Scandinavia",
    ]
    y = CY + Inches(0.48)
    for it in items_l:
        dot(sl, it, Inches(0.40), y, Inches(5.90))
        y += Inches(0.52)

    # Question box
    BX = Inches(6.65); BW = Inches(6.30)
    rect(sl, BX, CY, BW, Inches(3.10), NAVY)
    rect(sl, BX, CY, BW, Inches(0.06), GOLD)
    txb(sl, "Our Research Question",
        BX+Inches(0.18), CY+Inches(0.10), BW-Inches(0.30), Inches(0.35),
        size=12, color=GOLD)
    txb(sl,
        '"Do weather data improve electricity\n'
        'price forecasts for France —\n'
        'and does the answer change depending\n'
        'on the market conditions?"',
        BX+Inches(0.18), CY+Inches(0.48), BW-Inches(0.30), Inches(2.42),
        size=15, bold=True, color=WHITE)

    txb(sl, "Four specific questions:",
        BX, CY+Inches(3.25), BW, Inches(0.35), size=14, bold=True, color=NAVY)
    subs = [
        "SQ1  Can machine learning beat a simple baseline?",
        "SQ2  Does adding weather improve ML forecasts?",
        "SQ3  Does the weather effect depend on the market context?",
        "SQ4  Does a better forecast translate into real trading profits?",
    ]
    y2 = CY + Inches(3.68)
    for s in subs:
        dot(sl, s, BX, y2, BW, size=13)
        y2 += Inches(0.46)


# ════════════════════════════════════════════════════════════════
# SLIDE 4 — DATA
# ════════════════════════════════════════════════════════════════
def s04():
    sl = prs.slides.add_slide(BLANK)
    bg(sl)
    header(sl, "02 — Our Data", "64,224 hourly observations · May 2018 – April 2025 · 3 data sources")
    footer(sl, 4)

    sources = [
        ("Grid & Market Data\n(ENTSO-E)", NAVY, [
            "Hourly electricity prices (what we forecast)",
            "Demand forecasts published by RTE (French TSO)",
            "Power generation: gas, nuclear, hydro, solar, wind",
            "Cross-border electricity flows (DE, ES, GB, BE)",
        ]),
        ("Historical Weather\n(ERA5 — Copernicus)", BLUE, [
            "Temperature at 2m across France",
            "Solar radiation (sunshine hours proxy)",
            "Wind speed (wind power proxy)",
            "Precipitation",
        ]),
        ("Energy Commodity\nPrices", GREEN, [
            "Natural gas — TTF hub (EUR/MWh)",
            "Coal — API2 benchmark (EUR/t)",
            "Derived: Heating Degree Days (cold spells)\nbased on 17°C threshold",
        ]),
    ]

    BW = Inches(4.11); x0 = Inches(0.35)
    for i, (title, col, items) in enumerate(sources):
        x = x0 + i*(BW+Inches(0.13))
        rect(sl, x, CY, BW, Inches(0.60), col)
        txb(sl, title, x, CY, BW, Inches(0.60),
            size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        rect(sl, x, CY+Inches(0.60), BW, Inches(3.45), WHITE)
        y2 = CY + Inches(0.76)
        for it in items:
            dot(sl, it, x+Inches(0.10), y2, BW-Inches(0.10), size=12, dc=col)
            y2 += Inches(0.56)

    rect(sl, 0, Inches(5.80), SW, Inches(0.06), GOLD)
    ks = [("64,224", "hourly data points"), ("7 years", "2018 – 2025"),
          ("19 variables", "incl. 4 weather measures"), ("8,640 h", "held-out test set")]
    xk = Inches(0.40); KW = Inches(3.0)
    for val, lab in ks:
        kpi(sl, xk, Inches(5.96), KW, Inches(0.98), val, lab)
        xk += KW + Inches(0.38)


# ════════════════════════════════════════════════════════════════
# SLIDE 5 — MODELS
# ════════════════════════════════════════════════════════════════
def s05():
    sl = prs.slides.add_slide(BLANK)
    bg(sl)
    header(sl, "03 — Four Forecasting Models",
           "Key design: Models B and C are identical except for weather — that isolates its exact contribution")
    footer(sl, 5)

    models = [
        ("A", "Simple Baseline", "\"Same as last week\"",
         "Forecast = price at the\nsame hour, 7 days ago\n\nNo machine learning.",
         MGREY, "Lower bound"),
        ("B", "Random Forest", "Without weather",
         "500 decision trees\nvoting together.\n\nAll data except\nweather variables.",
         BLUE, "Tests ML value alone"),
        ("C", "Random Forest", "With weather  ★",
         "Same 500 trees,\nsame setup as B.\n\nAll data including\nweather variables.",
         NAVY, "Our main model"),
        ("D", "XGBoost", "With weather",
         "Gradient boosting\n(different ML technique).\n\nAll data including\nweather variables.",
         GREEN, "Alternative ML"),
    ]

    BW = Inches(3.05); x0 = Inches(0.28)
    for i, (letter, name, tag, desc, col, role) in enumerate(models):
        x = x0 + i*(BW+Inches(0.12))
        rect(sl, x, CY, BW, Inches(0.60), col)
        txb(sl, f"Model {letter}", x, CY, BW, Inches(0.30),
            size=10, color=WHITE, align=PP_ALIGN.CENTER)
        txb(sl, name, x, CY+Inches(0.30), BW, Inches(0.30),
            size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        rect(sl, x, CY+Inches(0.60), BW, Inches(0.28), LGREY)
        txb(sl, tag, x, CY+Inches(0.60), BW, Inches(0.28),
            size=11, bold=True, color=col, align=PP_ALIGN.CENTER)
        rect(sl, x, CY+Inches(0.88), BW, Inches(2.52), WHITE)
        txb(sl, desc, x+Inches(0.12), CY+Inches(1.00), BW-Inches(0.22), Inches(2.28),
            size=12, color=DGREY)
        rect(sl, x, CY+Inches(3.40), BW, Inches(0.30), col)
        txb(sl, role, x, CY+Inches(3.40), BW, Inches(0.30),
            size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    txb(sl, "B vs C  =  pure weather effect  (everything else identical)",
        Inches(2.50), CY+Inches(3.88), Inches(8.30), Inches(0.40),
        size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    txb(sl,
        "Train: 2018 – April 2024 (50,184 h)   |   Test: May 2024 – April 2025 (8,640 h)   |   Strict time order — no peeking at future data",
        Inches(0.30), CY+Inches(4.40), Inches(12.7), Inches(0.38),
        size=12, color=MGREY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 6 — RESULTS
# ════════════════════════════════════════════════════════════════
def s06():
    sl = prs.slides.add_slide(BLANK)
    bg(sl)
    header(sl, "04 — Forecasting Results: Normal Market (May 2024 – Apr 2025)",
           "MAE = average error in €/MWh  |  R² = share of price variation explained  |  Hit rate = % of correct directions")
    footer(sl, 6)

    rows = [
        ["Model",                   "Avg. Error (MAE)", "Fit (R²)", "Error % (sMAPE)", "Direction OK"],
        ["A — Baseline (last week)", "33.09 €/MWh",    "0.16",     "70.9%",           "59.6%"],
        ["B — ML without weather",   "16.95 €/MWh",    "0.78",     "44.2%",           "65.5%"],
        ["C — ML with weather  ★",   "16.94 €/MWh",    "0.78",     "44.2%",           "65.4%"],
        ["D — XGBoost",              "19.14 €/MWh",    "0.66",     "45.4%",           "65.8%"],
    ]
    cw = [Inches(3.20), Inches(2.00), Inches(1.60), Inches(2.00), Inches(2.00)]
    pptx_table(sl, rows, cw, x=Inches(0.55), y=CY,
               row_h=Inches(0.50), hfill=NAVY, hcolor=WHITE,
               alt=LGREY, body=DGREY, hsize=11, bsize=12,
               hilite=2, hifill=RGBColor(0xE0,0xEB,0xFF),
               col_align=[PP_ALIGN.LEFT]+[PP_ALIGN.CENTER]*4)

    Y_BOX = CY + Inches(2.75)
    finds = [
        (GREEN, "SQ1 ✓  ML is twice as accurate",
         "ML cuts the average error from 33 to 17 €/MWh.\nR² jumps from 0.16 to 0.78: we now explain 78% of price swings."),
        (RED,   "SQ2 — A surprise",
         "Models B and C are virtually identical.\nAdding weather data changes almost nothing."),
        (GOLD,  "The real question →",
         "Is weather truly useless?\nOr is it already hidden inside other variables?"),
    ]
    xb = Inches(0.35); BW2 = Inches(4.15)
    for col, title, text in finds:
        rect(sl, xb, Y_BOX, BW2, Inches(1.82), WHITE)
        rect(sl, xb, Y_BOX, Inches(0.08), Inches(1.82), col)
        txb(sl, title, xb+Inches(0.16), Y_BOX+Inches(0.08),
            BW2-Inches(0.22), Inches(0.38), size=13, bold=True, color=col)
        txb(sl, text, xb+Inches(0.16), Y_BOX+Inches(0.50),
            BW2-Inches(0.22), Inches(1.20), size=12, color=DGREY)
        xb += BW2 + Inches(0.18)


# ════════════════════════════════════════════════════════════════
# SLIDE 7 — DM TEST + IRH
# ════════════════════════════════════════════════════════════════
def s07():
    sl = prs.slides.add_slide(BLANK)
    bg(sl)
    header(sl, "04 — Is the Difference Real? Statistical Significance Test",
           "Diebold-Mariano test — asks: are two models genuinely different, or just lucky?  |  p-value < 0.05 = significant")
    footer(sl, 7)

    dm_rows = [
        ["Comparison",                    "What it tests",                  "DM stat.", "p-value",    "Verdict"],
        ["C vs B\n(weather vs no weather)","Does weather add anything?",     "−0.565",  "0.572",      "No difference"],
        ["C vs A\n(ML vs baseline)",       "Is ML better than last week?",   "−12.481", "< 0.001",    "Yes ✓✓✓"],
        ["B vs A\n(ML vs baseline)",       "ML without weather vs baseline?", "−12.456", "< 0.001",    "Yes ✓✓✓"],
        ["D vs C\n(XGBoost vs RF)",        "Is XGBoost better than RF?",     "+3.847",  "< 0.001",    "XGB higher err."],
    ]
    cw7 = [Inches(1.80), Inches(3.70), Inches(1.50), Inches(1.50), Inches(2.05)]
    pptx_table(sl, dm_rows, cw7, x=Inches(0.40), y=CY,
               row_h=Inches(0.50), hfill=NAVY, hcolor=WHITE,
               alt=LGREY, body=DGREY, hsize=11, bsize=12,
               hilite=0, hifill=RGBColor(0xFF,0xEC,0xEC),
               col_align=[PP_ALIGN.LEFT, PP_ALIGN.LEFT,
                          PP_ALIGN.CENTER, PP_ALIGN.CENTER, PP_ALIGN.CENTER])

    txb(sl, "→  Weather makes no statistically significant difference in a normal market (p = 0.572 >> 0.05)",
        Inches(0.40), CY+Inches(2.72), Inches(12.55), Inches(0.36),
        size=13, bold=True, color=RED, align=PP_ALIGN.CENTER)

    txb(sl, "Why? — The Information Redundancy Hypothesis",
        Inches(0.40), CY+Inches(3.18), Inches(12.55), Inches(0.36),
        size=16, bold=True, color=NAVY)

    LW = Inches(6.20); RX = Inches(0.40)+LW+Inches(0.12)
    RW = Inches(6.20); TY = CY+Inches(3.62); ROW = Inches(0.72)

    pairs = [
        ("Gas price (8.3% of model importance)",
         "Gas demand rises with cold weather → gas price already encodes temperature.",
         "RTE demand forecast (top variable)",
         "France's grid operator uses weather models internally → demand forecast IS a weather proxy."),
        ("Price lags (70% of importance)",
         "In a calm market, yesterday's price predicts tomorrow's better than weather.",
         "Raw weather data (ERA5, all 4 variables)",
         "Only 2.2% of importance — marginal once gas price and demand are in the model."),
    ]
    for row_i, (lt, ld, rt, rd) in enumerate(pairs):
        y = TY + row_i * ROW * 1.10
        rect(sl, Inches(0.40), y, Inches(0.08), ROW, BLUE)
        txb(sl, lt, Inches(0.58), y, LW-Inches(0.20), Inches(0.30),
            size=12, bold=True, color=NAVY)
        txb(sl, ld, Inches(0.58), y+Inches(0.30), LW-Inches(0.20), Inches(0.36),
            size=11, color=DGREY)
        rect(sl, RX, y, Inches(0.08), ROW, GOLD)
        txb(sl, rt, RX+Inches(0.18), y, RW-Inches(0.20), Inches(0.30),
            size=12, bold=True, color=NAVY)
        txb(sl, rd, RX+Inches(0.18), y+Inches(0.30), RW-Inches(0.20), Inches(0.36),
            size=11, color=DGREY)


# ════════════════════════════════════════════════════════════════
# SLIDE 8 — FEATURE IMPORTANCE
# ════════════════════════════════════════════════════════════════
def s08():
    sl = prs.slides.add_slide(BLANK)
    bg(sl)
    header(sl, "04 — What Does the Model Actually Use?",
           "Variable importance in the Random Forest (Model C) — how much each group contributes to accuracy")
    footer(sl, 8)

    txb(sl,
        ("The chart on the right shows what the model relies\n"
         "on most when predicting tomorrow's price:\n\n"
         "Price patterns (lags) dominate at 70%\n"
         "→ Electricity prices are mean-reverting: 'similar\n"
         "    hour last week' is a strong signal.\n\n"
         "Gas price contributes 8.3%, coal 6.0%\n"
         "→ These fuel costs drive marginal plant costs.\n"
         "    They already embed the weather signal.\n\n"
         "All 4 weather variables combined: only 2.2%\n"
         "→ Not because weather doesn't matter in reality —\n"
         "    but because it's already captured indirectly.\n\n"
         "Conclusion: weather is useful information,\n"
         "but it arrives encoded in other variables first."),
        Inches(0.40), CY, Inches(5.90), Inches(5.40), size=13, color=DGREY)

    CHART_X = Inches(6.55); CHART_Y = CY+Inches(0.30)
    MAX_BW  = Inches(4.30); BAR_H   = Inches(0.44); GAP = Inches(0.14)
    LABEL_W = Inches(1.80)

    txb(sl, "Contribution to forecast accuracy (%)",
        CHART_X, CHART_Y-Inches(0.38), Inches(6.20), Inches(0.34),
        size=12, bold=True, color=NAVY)

    features = [
        ("Price lags (7 variables)", 70.0, NAVY),
        ("TTF gas price",             8.3,  BLUE),
        ("Coal API2 price",           6.0,  BLUE),
        ("Demand forecast",           5.8,  BLUE),
        ("Cross-border flows",        4.5,  MGREY),
        ("Generation mix",            3.2,  MGREY),
        ("Weather — ERA5 (4 vars)",   2.2,  RED),
    ]
    for fi, (fname, pct, col) in enumerate(features):
        y = CHART_Y + fi*(BAR_H+GAP)
        bw = MAX_BW * pct / 70.0
        rect(sl, CHART_X, y, bw, BAR_H, col)
        if bw >= Inches(0.55):
            txb(sl, f"{pct}%", CHART_X, y, max(Inches(0.3), bw-Inches(0.06)), BAR_H,
                size=11, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
        lx = CHART_X + MAX_BW + Inches(0.10)
        txb(sl, fname, lx, y, LABEL_W, BAR_H, size=11, color=DGREY)

    rect(sl, CHART_X, CHART_Y+7*(BAR_H+GAP),
         MAX_BW+LABEL_W+Inches(0.10), Inches(0.05), GOLD)
    txb(sl, "In a normal market: focus on gas prices\nand demand data, not weather feeds.",
        CHART_X, CHART_Y+7*(BAR_H+GAP)+Inches(0.10),
        MAX_BW+LABEL_W, Inches(0.70), size=12, color=MGREY)


# ════════════════════════════════════════════════════════════════
# SLIDE 9 — CRISIS 2022
# ════════════════════════════════════════════════════════════════
def s09():
    sl = prs.slides.add_slide(BLANK)
    bg(sl)
    rect(sl, 0, 0, SW, HEADER_H, RGBColor(0xA0,0x20,0x20))
    rect(sl, 0, HEADER_H, SW, ACCENT_H, GOLD)
    txb(sl, "05 — The 2022 Energy Crisis: Weather Finally Matters",
        Inches(0.45), Inches(0.18), Inches(12.4), Inches(0.72),
        size=23, bold=True, color=WHITE)
    txb(sl, "Same models, same test method — applied to the crisis year",
        Inches(0.45), Inches(0.90), Inches(12.4), Inches(0.38), size=12, color=GOLD)
    footer(sl, 9)

    txb(sl, "We re-ran the test on 2022 — and got the opposite result:",
        Inches(0.40), CY, Inches(12.55), Inches(0.38),
        size=16, bold=True, color=DGREY)

    rows9 = [
        ["What we measured",        "Normal market (2024-25)", "Crisis year (2022)",  "Meaning"],
        ["Average error (MAE)",     "17 €/MWh",               "67 €/MWh",           "3.9× harder to forecast"],
        ["Model fit (R²)",          "0.78  (78% explained)",  "0.48  (48%)",         "Much less predictable"],
        ["Weather adds value?",     "No  (p = 0.572)",        "YES  (p < 0.001)",    "Complete reversal"],
        ["DM statistic (C vs B)",   "−0.565  →  n.s.",        "−13.27  →  ✓✓✓",     "Statistically certain"],
    ]
    cw9 = [Inches(2.80), Inches(2.60), Inches(2.60), Inches(3.00)]
    pptx_table(sl, rows9, cw9, x=Inches(0.40), y=CY+Inches(0.48),
               row_h=Inches(0.48), hfill=RGBColor(0xA0,0x20,0x20), hcolor=WHITE,
               alt=LGREY, body=DGREY, hsize=11, bsize=12,
               col_align=[PP_ALIGN.LEFT, PP_ALIGN.CENTER,
                          PP_ALIGN.CENTER, PP_ALIGN.LEFT])

    txb(sl, "Why did the result flip?",
        Inches(0.40), CY+Inches(3.52), Inches(12.55), Inches(0.36),
        size=15, bold=True, color=NAVY)
    reasons = [
        "Gas prices went haywire (Russia war) — they no longer tracked weather demand, breaking the proxy.",
        "30 GW of French nuclear capacity was offline at once (corrosion crisis) — another proxy lost.",
        "With both 'weather encoders' broken, the model needed direct weather data to fill the gap.",
        "Result: adding weather cut the error by ≈15% in 2022 — both statistically and economically significant.",
    ]
    y = CY + Inches(3.98)
    for r in reasons:
        dot(sl, r, Inches(0.40), y, Inches(12.55), size=12, dc=RED)
        y += Inches(0.44)


# ════════════════════════════════════════════════════════════════
# SLIDE 10 — REGIME DEPENDENCE
# ════════════════════════════════════════════════════════════════
def s10():
    sl = prs.slides.add_slide(BLANK)
    bg(sl)
    header(sl, "05 — Our Core Finding: Weather Value Depends on Market Context",
           "The Information Redundancy Hypothesis — weather matters only when its proxies break down")
    footer(sl, 10)

    COL_W = Inches(5.90); DIV_W = Inches(0.60)
    LX = Inches(0.35); DIV_X = LX+COL_W+Inches(0.08)
    RX = DIV_X+DIV_W+Inches(0.08); COL_H = Inches(4.78)

    rect(sl, LX, CY, COL_W, COL_H, WHITE)
    rect(sl, LX, CY, COL_W, Inches(0.06), GREEN)
    txb(sl, "NORMAL MARKET  (2024-25)",
        LX+Inches(0.12), CY+Inches(0.10), COL_W-Inches(0.20), Inches(0.36),
        size=14, bold=True, color=GREEN)

    stable = [
        ("Gas price tracks cold weather demand", False, DGREY),
        ("↳  Gas price = indirect weather signal", True,  GREEN),
        ("Demand forecasts built with weather models", False, DGREY),
        ("↳  Demand forecast = another weather proxy", True, GREEN),
        ("Nuclear output is stable and predictable", False, DGREY),
        ("", False, DGREY),
        ("Direct weather data → REDUNDANT", True,  GREEN),
        ("Statistical test: p = 0.572  (no difference)", True,  DGREY),
        ("MAE improvement from weather: < 0.1%", False, DGREY),
    ]
    y = CY+Inches(0.58)
    for txt, bold, col in stable:
        if txt:
            txb(sl, txt, LX+Inches(0.14), y, COL_W-Inches(0.22), Inches(0.40),
                size=12, bold=bold, color=col)
        y += Inches(0.42)

    rect(sl, DIV_X, CY, DIV_W, COL_H, LGREY)
    txb(sl, "REGIME\nSHIFT", DIV_X, CY+Inches(2.0), DIV_W, Inches(0.80),
        size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    rect(sl, RX, CY, COL_W, COL_H, WHITE)
    rect(sl, RX, CY, COL_W, Inches(0.06), RED)
    txb(sl, "CRISIS MARKET  (2022)",
        RX+Inches(0.12), CY+Inches(0.10), COL_W-Inches(0.20), Inches(0.36),
        size=14, bold=True, color=RED)

    crisis = [
        ("Gas prices disconnected from cold weather", False, DGREY),
        ("↳  Gas price proxy BROKEN", True,  RED),
        ("30 GW nuclear offline (corrosion crisis)", False, DGREY),
        ("↳  Demand forecast less reliable", True,  RED),
        ("Price patterns become chaotic", False, DGREY),
        ("", False, DGREY),
        ("Direct weather data → ESSENTIAL", True,  RED),
        ("Statistical test: p < 0.001  (clear difference)", True,  DGREY),
        ("DM statistic: −13.27", False, DGREY),
    ]
    y = CY+Inches(0.58)
    for txt, bold, col in crisis:
        if txt:
            txb(sl, txt, RX+Inches(0.14), y, COL_W-Inches(0.22), Inches(0.40),
                size=12, bold=bold, color=col)
        y += Inches(0.42)

    rect(sl, Inches(0.35), CY+COL_H+Inches(0.12), SW-Inches(0.70), Inches(0.64), NAVY)
    txb(sl,
        "THESIS CONTRIBUTION: Weather data should be treated as a RISK MANAGEMENT signal,"
        " not a routine forecasting input. It matters most precisely when markets break down.",
        Inches(0.55), CY+COL_H+Inches(0.16), SW-Inches(1.10), Inches(0.56),
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 11 — TRADING STRATEGY
# ════════════════════════════════════════════════════════════════
def s11():
    sl = prs.slides.add_slide(BLANK)
    bg(sl)
    header(sl, "06 — Trading Backtest: Turning Forecasts into Trades",
           "EPEX SPOT France day-ahead  |  1 MW virtual position  |  May 2024 – April 2025")
    footer(sl, 11)

    BW = Inches(6.12)
    rect(sl, Inches(0.35), CY, BW, Inches(2.35), WHITE)
    rect(sl, Inches(0.35), CY, Inches(0.08), Inches(2.35), NAVY)
    txb(sl, "Trading Signal",
        Inches(0.55), CY+Inches(0.08), BW-Inches(0.24), Inches(0.36),
        size=14, bold=True, color=NAVY)
    txb(sl,
        ("Rule: if our forecast for tomorrow (hour h) is\n"
         "higher than today's actual price → BUY  (+1)\n"
         "lower than today's actual price → SELL (−1)\n\n"
         "signal(h) = sign[ forecast(h, tomorrow) − actual(h, today) ]\n\n"
         "Safety filter: if predicted move < 2 €/MWh → no trade\n"
         "(transaction costs would exceed expected profit)"),
        Inches(0.55), CY+Inches(0.50), BW-Inches(0.24), Inches(1.76),
        size=12, color=DGREY)

    RX = Inches(0.35)+BW+Inches(0.14); RW = SW-RX-Inches(0.35)
    rect(sl, RX, CY, RW, Inches(2.35), WHITE)
    rect(sl, RX, CY, Inches(0.08), Inches(2.35), GOLD)
    txb(sl, "Profit & Loss",
        RX+Inches(0.18), CY+Inches(0.08), RW-Inches(0.24), Inches(0.36),
        size=14, bold=True, color=NAVY)
    txb(sl,
        ("Gross profit(h) = signal × [actual(h, tomorrow) − actual(h, today)]\n\n"
         "Net profit(h) = Gross − transaction cost × |signal|\n\n"
         "Position: 1 MW  ·  No borrowing or leverage\n"
         "Units: euros per hour  (1 MW × 1 h × €/MWh)"),
        RX+Inches(0.18), CY+Inches(0.50), RW-Inches(0.24), Inches(1.76),
        size=12, color=DGREY)

    txb(sl, "Transaction Cost Scenarios (EPEX SPOT official fee schedule)",
        Inches(0.35), CY+Inches(2.52), Inches(12.55), Inches(0.36),
        size=14, bold=True, color=NAVY)
    scen = [
        ("Optimistic",  "0.10 €/MWh", "Exchange fee only. Large participant.", GREEN),
        ("Central",     "0.30 €/MWh", "Fee + bid-ask spread. Typical trader.", NAVY),
        ("Pessimistic", "0.60 €/MWh", "Central + settlement risk. Upper bound.", RED),
    ]
    xs = Inches(0.35); SW3 = (SW-Inches(0.70))/3
    for name, cost, desc, col in scen:
        rect(sl, xs, CY+Inches(2.96), SW3, Inches(0.40), col)
        txb(sl, f"{name}  —  {cost}", xs+Inches(0.10), CY+Inches(2.96),
            SW3-Inches(0.15), Inches(0.40), size=12, bold=True, color=WHITE)
        rect(sl, xs, CY+Inches(3.36), SW3, Inches(1.18), WHITE)
        txb(sl, desc, xs+Inches(0.10), CY+Inches(3.44),
            SW3-Inches(0.15), Inches(1.05), size=12, color=DGREY)
        xs += SW3

    txb(sl,
        ("Risk metrics: Sharpe ratio = (avg daily P&L ÷ daily volatility) × √365  "
         "|  Max Drawdown = biggest loss streak in €/MW  "
         "|  Calmar = yearly profit ÷ max drawdown"),
        Inches(0.35), CY+Inches(4.72), Inches(12.55), Inches(0.44),
        size=11, color=MGREY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 12 — BACKTEST RESULTS
# ════════════════════════════════════════════════════════════════
def s12():
    sl = prs.slides.add_slide(BLANK)
    bg(sl)
    header(sl, "06 — Trading Results: Central Cost Scenario (0.30 €/MWh)",
           "12-month test · 1 MW position · Net of transaction costs")
    footer(sl, 12)

    rows12 = [
        ["Model",              "Net Profit (€/MW/yr)", "Sharpe ratio", "Worst loss (€/MW)", "Calmar", "Win rate", "Active"],
        ["A — Baseline",       "85,482 €",              "10.5",         "2,367 €",            "37",     "62%",     "95%"],
        ["B — ML w/o weather", "136,884 €",             "19.4",         "433 €",              "321",    "75%",     "88%"],
        ["C — ML with weather","136,711 €",             "19.5",         "422 €",              "328",    "75%",     "88%"],
        ["D — XGBoost",        "142,012 €",             "18.8",         "1,351 €",            "107",    "75%",     "93%"],
        ["Buy-and-hold ref.",  "−76 €",                 "≈ 0",          "4,521 €",            "—",      "48%",     "100%"],
    ]
    cw12 = [Inches(2.30), Inches(1.90), Inches(1.30), Inches(1.80),
            Inches(1.08), Inches(0.95), Inches(0.90)]
    pptx_table(sl, rows12, cw12, x=Inches(0.15), y=CY,
               row_h=Inches(0.44), hfill=NAVY, hcolor=WHITE,
               alt=LGREY, body=DGREY, hsize=10, bsize=11,
               hilite=2, hifill=RGBColor(0xE0,0xEB,0xFF),
               col_align=[PP_ALIGN.LEFT]+[PP_ALIGN.CENTER]*6)

    Y_TK = CY+Inches(3.10); BW_TK = Inches(3.18)
    takes = [
        (GREEN, "ML earns 60% more",
         "136K€ vs 85K€ per year for 1 MW.\nBuy-and-hold earns nothing → it's pure skill."),
        (NAVY,  "RF has far less risk",
         "Worst drawdown: RF = 422€, XGBoost = 1,351€.\nCalmar: RF = 328 vs XGBoost = 107."),
        (BLUE,  "Weather still neutral",
         "136,884€ vs 136,711€ — nearly identical.\nEconomic confirmation of the statistical test."),
        (GOLD,  "Robust across cost levels",
         "Sharpe ratio changes by only 0.5 points\nfrom cheapest to most expensive cost scenario."),
    ]
    xb = Inches(0.15)
    for col, title, text in takes:
        rect(sl, xb, Y_TK, BW_TK, Inches(1.85), WHITE)
        rect(sl, xb, Y_TK, Inches(0.08), Inches(1.85), col)
        txb(sl, title, xb+Inches(0.16), Y_TK+Inches(0.08),
            BW_TK-Inches(0.22), Inches(0.36), size=12, bold=True, color=col)
        txb(sl, text, xb+Inches(0.16), Y_TK+Inches(0.48),
            BW_TK-Inches(0.22), Inches(1.20), size=11, color=DGREY)
        xb += BW_TK+Inches(0.09)


# ════════════════════════════════════════════════════════════════
# SLIDE 13 — ROBUSTNESS
# ════════════════════════════════════════════════════════════════
def s13():
    sl = prs.slides.add_slide(BLANK)
    bg(sl)
    header(sl, "06 — Is It Consistent? Robustness Checks",
           "Monthly performance  |  Cost sensitivity  |  12 months with no cherry-picking")
    footer(sl, 13)

    txb(sl, "Monthly Sharpe ratio — each calendar month, independently",
        Inches(0.35), CY, Inches(8.20), Inches(0.35), size=14, bold=True, color=NAVY)

    MONTHS  = ["May","Jun","Jul","Aug","Sep","Oct","Nov","Dec","Jan","Feb","Mar","Apr"]
    NAIVE_S = [9.74,11.73,17.88,7.92,17.15,15.41,15.92,9.98,5.89,3.91,8.21,12.38]
    RF_S    = [19.31,23.83,28.28,16.52,27.24,22.93,17.90,21.89,17.66,21.74,17.59,16.80]
    XGB_S   = [25.48,22.55,29.75,16.99,33.97,25.94,20.31,26.27,20.32,21.87,22.91,-0.99]

    CX=Inches(0.35); CY2=CY+Inches(0.42); CW=Inches(8.20); CH=Inches(3.30)
    MAX_S=36.0; N=len(MONTHS)
    GW=CW/N; BW2=GW*0.28

    rect(sl, CX, CY2, CW, CH, WHITE)
    for mi, month in enumerate(MONTHS):
        gx = CX+mi*GW
        h_n = max(Inches(0.01), CH*NAIVE_S[mi]/MAX_S)
        rect(sl, gx+Inches(0.01), CY2+CH-h_n, BW2, h_n, MGREY)
        h_r = max(Inches(0.01), CH*RF_S[mi]/MAX_S)
        rect(sl, gx+BW2+Inches(0.01), CY2+CH-h_r, BW2, h_r, NAVY)
        h_x = max(Inches(0.01), CH*abs(XGB_S[mi])/MAX_S)
        rect(sl, gx+2*BW2+Inches(0.01), CY2+CH-h_x, BW2, h_x,
             RED if XGB_S[mi]<0 else GREEN)
        txb(sl, month, gx, CY2+CH+Inches(0.02), GW, Inches(0.28),
            size=8, color=MGREY, align=PP_ALIGN.CENTER)

    LY = CY2+CH+Inches(0.38); LX = Inches(0.35)
    for col, lbl in [(MGREY,"Baseline"),(NAVY,"RF + weather"),(GREEN,"XGBoost"),(RED,"XGB (negative)")]:
        rect(sl, LX, LY, Inches(0.24), Inches(0.20), col)
        txb(sl, lbl, LX+Inches(0.30), LY-Inches(0.02), Inches(1.40), Inches(0.26),
            size=11, color=DGREY)
        LX += Inches(2.00)

    txb(sl,
        ("RF models: positive Sharpe every single month — 12 out of 12.\n"
         "Range: 16.5 to 28.3. No single month drives the result.\n"
         "XGBoost: negative Sharpe in April 2025 — less stable across regimes."),
        Inches(0.35), LY+Inches(0.30), Inches(8.20), Inches(0.90), size=12, color=DGREY)

    txb(sl, "Does the cost assumption change anything?",
        Inches(8.80), CY, Inches(4.20), Inches(0.35),
        size=14, bold=True, color=NAVY)
    rows_cs = [
        ["Model",        "0.10 €/MWh", "0.30 €/MWh", "0.60 €/MWh"],
        ["Baseline",     "10.73",       "10.53",       "10.23"],
        ["RF w/o wx",    "19.64",       "19.44",       "19.14"],
        ["RF with wx ★", "19.66",       "19.46",       "19.16"],
        ["XGBoost",      "18.98",       "18.78",       "18.47"],
    ]
    cw_cs = [Inches(1.40), Inches(0.92), Inches(0.92), Inches(1.08)]
    pptx_table(sl, rows_cs, cw_cs, x=Inches(8.80), y=CY+Inches(0.42),
               row_h=Inches(0.50), hfill=NAVY, hcolor=WHITE,
               alt=LGREY, body=DGREY, hsize=10, bsize=11,
               hilite=2, hifill=RGBColor(0xE0,0xEB,0xFF))
    txb(sl, "Maximum change:\n≈ 0.5 Sharpe points\nacross all 3 scenarios.\nNot sensitive to costs.",
        Inches(8.80), CY+Inches(3.10), Inches(4.20), Inches(1.20), size=12, color=DGREY)


# ════════════════════════════════════════════════════════════════
# SLIDE 14 — LIMITATIONS
# ════════════════════════════════════════════════════════════════
def s14():
    sl = prs.slides.add_slide(BLANK)
    bg(sl)
    header(sl, "07 — What We'd Do Differently — Honest Limitations",
           "Self-critique is part of rigorous research")
    footer(sl, 14)

    limits = [
        ("Only 12 months tested",
         "Our normal-market test covers one year only. A full rolling walk-forward — retrain every year, test the next — would be more convincing."),
        ("Perfect trade execution assumed",
         "We assume bids are always matched at the market price. In reality, large or systematic trades can shift the price against you."),
        ("Fixed position size",
         "We always trade 1 MW, regardless of how confident the model is. Sizing up on high-confidence signals could improve returns further."),
        ("CO₂ price data incomplete",
         "We could not source a reliable EU Allowance (EUA) series. 'Clean spark spreads' — a key market signal — were therefore approximated."),
        ("Sharpe ratios look very high",
         "19× annualised sounds extreme vs equities. But electricity is a different asset: no leverage, prices mean-revert strongly, market is less competitive."),
        ("No uncertainty quantification",
         "Our model gives a point forecast. A probability range (e.g. 'price between 50 and 80 €/MWh') would improve the no-trade filter calibration."),
    ]

    y = CY+Inches(0.05); ROW = Inches(0.84)
    for title, text in limits:
        rect(sl, Inches(0.35), y, Inches(0.08), ROW-Inches(0.08), GOLD)
        txb(sl, title, Inches(0.55), y, Inches(3.50), Inches(0.34),
            size=13, bold=True, color=NAVY)
        txb(sl, text, Inches(4.20), y, Inches(8.90), ROW-Inches(0.08),
            size=12, color=DGREY)
        y += ROW


# ════════════════════════════════════════════════════════════════
# SLIDE 15 — CONCLUSIONS
# ════════════════════════════════════════════════════════════════
def s15():
    sl = prs.slides.add_slide(BLANK)
    bg(sl)
    header(sl, "07 — Conclusions"); footer(sl, 15)

    conclusions = [
        ("SQ1 ✓",  GREEN, "Machine learning beats the simple baseline by a wide margin",
         "Average error drops from 33 to 17 €/MWh (−49%). The model explains 78% of price variation. Trading Sharpe: 10.5 → 19.5."),
        ("SQ2 ✗/✓", GOLD, "Weather data: neutral in calm markets, critical in crises",
         "Normal market: adding weather changes nothing (p = 0.572). 2022 crisis: weather becomes essential (DM = −13.27, p < 0.001)."),
        ("SQ3 ✓",  NAVY, "The Information Redundancy Hypothesis is confirmed",
         "In calm markets, gas prices and demand forecasts already encode weather. The direct signal only matters when those proxies break."),
        ("SQ4 ✓",  BLUE, "Better forecasts translate directly into trading profits",
         "RF earns 136K €/MW/yr, baseline earns 85K. Buy-and-hold earns zero → profits come from genuine directional skill."),
    ]

    y = CY+Inches(0.04); ROW = Inches(1.38)
    for label, col, title, text in conclusions:
        rect(sl, Inches(0.35), y, Inches(0.10), ROW, col)
        rect(sl, Inches(0.55), y+Inches(0.04), Inches(1.20), Inches(0.38), col)
        txb(sl, label, Inches(0.55), y+Inches(0.04), Inches(1.20), Inches(0.38),
            size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txb(sl, title, Inches(1.90), y+Inches(0.04), Inches(11.0), Inches(0.38),
            size=15, bold=True, color=col)
        txb(sl, text, Inches(1.90), y+Inches(0.48), Inches(11.0), Inches(0.82),
            size=12, color=DGREY)
        y += ROW

    rect(sl, Inches(0.35), y+Inches(0.06), SW-Inches(0.70), Inches(0.60), NAVY)
    txb(sl,
        "Key message: treat weather as a risk management signal — "
        "monitor it continuously, activate it aggressively when markets break down.",
        Inches(0.55), y+Inches(0.10), SW-Inches(1.10), Inches(0.52),
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 16 — Q&A
# ════════════════════════════════════════════════════════════════
def s16():
    sl = prs.slides.add_slide(BLANK)
    bg(sl)
    header(sl, "Anticipated Questions — Talking Points",
           "20-40 min discussion with the thesis director")
    footer(sl, 16)

    qas = [
        ("Why sMAPE and not the standard MAPE error metric?",
         "MAPE is meaningless near zero prices and explodes with negative prices — both occur in France. sMAPE stays between 0% and 200% and treats over/under-forecasts symmetrically."),
        ("Is the Diebold-Mariano test valid on 8,640 data points?",
         "Yes. The Harvey-Leybourne-Newbold correction is specifically designed for small samples. With n = 8,640, the test has very high statistical power."),
        ("Why use last week (lag 168h) and not yesterday (lag 24h) as baseline?",
         "Electricity prices have a strong day-of-week pattern (Monday peak, Sunday trough). Last week same hour is the industry-standard baseline (Weron, 2014)."),
        ("Why does XGBoost earn more but rank lower than RF?",
         "XGBoost fits the training data more aggressively and then collapses in April 2025 (Sharpe = −0.99). RF is more stable. On risk-adjusted return (Calmar 328 vs 107), RF wins."),
        ("A Sharpe ratio of 19 seems unrealistically high — is the backtest flawed?",
         "No — electricity markets are structurally different from equities: no leverage, strongly mean-reverting prices, less sophisticated competition. Compare models to each other, not to hedge fund standards."),
    ]

    y = CY+Inches(0.04)
    for q, a in qas:
        rect(sl, Inches(0.35), y, SW-Inches(0.70), Inches(0.30), NAVY)
        txb(sl, f"Q: {q}", Inches(0.50), y+Inches(0.02), SW-Inches(1.0), Inches(0.28),
            size=12, bold=True, color=WHITE)
        rect(sl, Inches(0.35), y+Inches(0.30), SW-Inches(0.70), Inches(0.64), WHITE)
        txb(sl, a, Inches(0.50), y+Inches(0.34), SW-Inches(1.0), Inches(0.58),
            size=11, color=DGREY)
        y += Inches(1.02)


# ════════════════════════════════════════════════════════════════
# SLIDE 17 — THANK YOU
# ════════════════════════════════════════════════════════════════
def s17():
    sl = prs.slides.add_slide(BLANK)
    rect(sl, 0, 0, SW, SH, NAVY)
    rect(sl, 0, 0, Inches(0.22), SH, GOLD)
    rect(sl, Inches(0.22), Inches(3.62), SW, Inches(0.06), GOLD)
    txb(sl, "Thank you",
        Inches(0.55), Inches(1.40), Inches(12.0), Inches(1.00),
        size=52, bold=True, color=WHITE)
    txb(sl, "Questions & Discussion",
        Inches(0.55), Inches(2.52), Inches(12.0), Inches(0.55), size=22, color=GOLD)
    txb(sl, "Lyam Oumedjeber   lyam.oumedjeber@edhec.com",
        Inches(0.55), Inches(3.88), Inches(11.0), Inches(0.38), size=14, color=WHITE)
    txb(sl, "Leo Cambreleng   leo.cambreleng@edhec.com",
        Inches(0.55), Inches(4.32), Inches(11.0), Inches(0.38), size=14, color=WHITE)
    txb(sl, "Thesis Director: Prof. Milos Vulanovic",
        Inches(0.55), Inches(4.84), Inches(11.0), Inches(0.38), size=13, color=MGREY)
    txb(sl,
        ("Core finding: weather data is neutral in normal markets (p = 0.572)\n"
         "but decisive in crises (DM = −13.27, p < 0.001).\n"
         "→ Information Redundancy Hypothesis — regime-conditional weather contribution."),
        Inches(0.55), Inches(5.55), Inches(11.0), Inches(0.90), size=12, color=GOLD)
    txb(sl, "EDHEC Business School  |  MSc DAAI  |  2025-2026",
        Inches(0.55), Inches(6.60), Inches(12.0), Inches(0.38), size=11, color=MGREY)


# ════════════════════════════════════════════════════════════════
# BUILD
# ════════════════════════════════════════════════════════════════
print("Building slides...")
s01(); print("  1/17")
s02(); print("  2/17")
s03(); print("  3/17")
s04(); print("  4/17")
s05(); print("  5/17")
s06(); print("  6/17")
s07(); print("  7/17")
s08(); print("  8/17")
s09(); print("  9/17")
s10(); print(" 10/17")
s11(); print(" 11/17")
s12(); print(" 12/17")
s13(); print(" 13/17")
s14(); print(" 14/17")
s15(); print(" 15/17")
s16(); print(" 16/17")
s17(); print(" 17/17")

prs.save(str(OUT))
print(f"\nSaved → {OUT}")
print(f"Size: {OUT.stat().st_size/1024:.0f} KB")
